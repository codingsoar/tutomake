"""
Document Exporter Module
Exports tutorials to PDF, PowerPoint, Markdown, PNG sequence formats
"""
import os
from typing import Callable, Optional
from ..model import Tutorial, Step


class DocumentExporter:
    """Export tutorial to document formats."""
    
    def __init__(self, tutorial: Tutorial, progress_callback: Optional[Callable[[int], None]] = None):
        self.tutorial = tutorial
        self.progress_callback = progress_callback
    
    def export_png_sequence(self, output_dir: str) -> bool:
        """Export each step as a PNG image with hitbox overlay."""
        import cv2
        import numpy as np
        
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
        
        total_steps = len(self.tutorial.steps)
        
        for i, step in enumerate(self.tutorial.steps):
            # Load image from step or video frame
            if step.image_path and os.path.exists(step.image_path):
                img = cv2.imread(step.image_path)
            elif self.tutorial.video_path and os.path.exists(self.tutorial.video_path):
                cap = cv2.VideoCapture(self.tutorial.video_path)
                fps = cap.get(cv2.CAP_PROP_FPS) or 24.0
                frame_num = int(step.timestamp * fps)
                cap.set(cv2.CAP_PROP_POS_FRAMES, frame_num)
                ret, img = cap.read()
                cap.release()
                if not ret:
                    continue
            else:
                continue
            
            # Draw hitbox overlay
            img = self._draw_overlay(img, step, i + 1)
            
            # Save
            output_path = os.path.join(output_dir, f"step_{i+1:03d}.png")
            cv2.imwrite(output_path, img)
            
            if self.progress_callback:
                self.progress_callback(int((i + 1) / total_steps * 100))
        
        print(f"Exported PNG sequence to: {output_dir}")
        return True
    
    def _draw_overlay(self, img, step: Step, step_num: int):
        """Draw hitbox and step number overlay."""
        import cv2
        
        def hex_to_bgr(hex_color):
            """Convert hex color (#RRGGBB or #RRGGBBAA) to BGR tuple."""
            hex_color = hex_color.lstrip('#')
            if len(hex_color) >= 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return (b, g, r)
            return (0, 0, 255)  # Default red
        
        if step.action_type == "keyboard":
            x, y = step.x, step.y
            cv2.rectangle(img, (x, y), (x + 300, y + 50), (255, 150, 0), 2)
            cv2.putText(img, step.keyboard_input, (x + 10, y + 35),
                       cv2.FONT_HERSHEY_SIMPLEX, 0.8, (255, 255, 255), 2)
        else:
            x, y, w, h = step.x, step.y, step.width, step.height
            overlay = img.copy()
            
            # Get colors from step's hitbox styling (same defaults as Player)
            line_color = hex_to_bgr(step.hitbox_line_color) if step.hitbox_line_color else hex_to_bgr("#FF0000")
            line_width = step.hitbox_line_width if step.hitbox_line_width else 2
            
            # Get fill color and opacity
            fill_color_str = step.hitbox_fill_color or "#FF0000"
            fill_color = hex_to_bgr(fill_color_str)
            fill_opacity = step.hitbox_fill_opacity if hasattr(step, 'hitbox_fill_opacity') else 20
            fill_alpha = fill_opacity / 100.0
            
            # Get line style
            line_style = step.hitbox_line_style if step.hitbox_line_style else "solid"
            
            if step.shape == "circle":
                center = (x + w // 2, y + h // 2)
                radius = max(w, h) // 2
                cv2.circle(overlay, center, radius, fill_color, -1)
                img = cv2.addWeighted(overlay, fill_alpha, img, 1 - fill_alpha, 0)
                
                # Draw circle with line style
                if line_style == "dashed":
                    # Draw dashed circle using multiple arcs
                    for angle in range(0, 360, 20):
                        start_angle = angle
                        end_angle = angle + 10
                        cv2.ellipse(img, center, (radius, radius), 0, start_angle, end_angle, line_color, line_width)
                elif line_style == "dotted":
                    # Draw dotted circle using small arcs
                    for angle in range(0, 360, 10):
                        start_angle = angle
                        end_angle = angle + 3
                        cv2.ellipse(img, center, (radius, radius), 0, start_angle, end_angle, line_color, line_width)
                else:
                    cv2.circle(img, center, radius, line_color, line_width)
            else:
                cv2.rectangle(overlay, (x, y), (x + w, y + h), fill_color, -1)
                img = cv2.addWeighted(overlay, fill_alpha, img, 1 - fill_alpha, 0)
                
                # Draw rectangle with line style
                if line_style == "dashed":
                    # Draw dashed rectangle
                    dash_length = 15
                    gap_length = 10
                    # Top edge
                    for px in range(x, x + w, dash_length + gap_length):
                        cv2.line(img, (px, y), (min(px + dash_length, x + w), y), line_color, line_width)
                    # Bottom edge
                    for px in range(x, x + w, dash_length + gap_length):
                        cv2.line(img, (px, y + h), (min(px + dash_length, x + w), y + h), line_color, line_width)
                    # Left edge
                    for py in range(y, y + h, dash_length + gap_length):
                        cv2.line(img, (x, py), (x, min(py + dash_length, y + h)), line_color, line_width)
                    # Right edge
                    for py in range(y, y + h, dash_length + gap_length):
                        cv2.line(img, (x + w, py), (x + w, min(py + dash_length, y + h)), line_color, line_width)
                elif line_style == "dotted":
                    # Draw dotted rectangle
                    dot_spacing = 8
                    # Top edge
                    for px in range(x, x + w, dot_spacing):
                        cv2.circle(img, (px, y), line_width, line_color, -1)
                    # Bottom edge
                    for px in range(x, x + w, dot_spacing):
                        cv2.circle(img, (px, y + h), line_width, line_color, -1)
                    # Left edge
                    for py in range(y, y + h, dot_spacing):
                        cv2.circle(img, (x, py), line_width, line_color, -1)
                    # Right edge
                    for py in range(y, y + h, dot_spacing):
                        cv2.circle(img, (x + w, py), line_width, line_color, -1)
                else:
                    cv2.rectangle(img, (x, y), (x + w, y + h), line_color, line_width)
        
        # Step number badge
        cv2.circle(img, (50, 50), 30, (255, 100, 0), -1)
        cv2.putText(img, str(step_num), (35, 60), cv2.FONT_HERSHEY_SIMPLEX, 1, (255, 255, 255), 2)
        
        return img
    
    def export_markdown(self, output_path: str, image_dir: str = None) -> bool:
        """Export as Markdown document."""
        if image_dir:
            self.export_png_sequence(image_dir)
        
        lines = [
            f"# {self.tutorial.title}",
            "",
            "---",
            ""
        ]
        
        for i, step in enumerate(self.tutorial.steps):
            lines.append(f"## Step {i + 1}: {step.description}")
            lines.append("")
            
            if step.instruction:
                lines.append(f"> 💡 {step.instruction}")
                lines.append("")
            
            if image_dir:
                img_path = f"step_{i+1:03d}.png"
                lines.append(f"![Step {i+1}]({image_dir}/{img_path})")
                lines.append("")
            
            if step.action_type == "keyboard":
                lines.append(f"**Type:** `{step.keyboard_input}`")
            else:
                lines.append(f"**Action:** Click at ({step.x}, {step.y})")
            
            lines.append("")
            lines.append("---")
            lines.append("")
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        
        print(f"Exported Markdown to: {output_path}")
        return True
    
    def export_pdf(self, output_path: str) -> bool:
        """Export as PDF document (requires fpdf2)."""
        try:
            from fpdf import FPDF
        except ImportError:
            print("fpdf2 is required for PDF export. Install with: pip install fpdf2")
            return False
        
        import cv2
        import tempfile
        
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        
        total_steps = len(self.tutorial.steps)
        
        # Title page
        pdf.add_page()
        pdf.set_font('Arial', 'B', 24)
        pdf.cell(0, 60, self.tutorial.title, ln=True, align='C')
        pdf.set_font('Arial', '', 12)
        pdf.cell(0, 10, f"Total Steps: {total_steps}", ln=True, align='C')
        
        for i, step in enumerate(self.tutorial.steps):
            pdf.add_page()
            
            # Step header
            pdf.set_font('Arial', 'B', 16)
            pdf.cell(0, 10, f"Step {i + 1}: {step.description}", ln=True)
            
            if step.instruction:
                pdf.set_font('Arial', 'I', 11)
                pdf.multi_cell(0, 8, step.instruction)
                pdf.set_font('Arial', 'B', 16)
            
            pdf.ln(5)
            
            # Get image
            img = None
            if step.image_path and os.path.exists(step.image_path):
                img = cv2.imread(step.image_path)
            elif self.tutorial.video_path and os.path.exists(self.tutorial.video_path):
                cap = cv2.VideoCapture(self.tutorial.video_path)
                fps = cap.get(cv2.CAP_PROP_FPS) or 24.0
                cap.set(cv2.CAP_PROP_POS_FRAMES, int(step.timestamp * fps))
                ret, img = cap.read()
                cap.release()
            
            if img is not None:
                img = self._draw_overlay(img, step, i + 1)
                
                # Save temp image
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    cv2.imwrite(tmp.name, img)
                    pdf.image(tmp.name, x=10, w=190)
                    os.unlink(tmp.name)
            
            # Action description
            pdf.ln(5)
            pdf.set_font('Arial', '', 12)
            if step.action_type == "keyboard":
                pdf.cell(0, 10, f"Action: Type '{step.keyboard_input}'", ln=True)
            else:
                pdf.cell(0, 10, f"Action: Click at position ({step.x}, {step.y})", ln=True)
            
            if self.progress_callback:
                self.progress_callback(int((i + 1) / total_steps * 100))
        
        pdf.output(output_path)
        print(f"Exported PDF to: {output_path}")
        return True
    
    def export_pptx(self, output_path: str) -> bool:
        """Export as PowerPoint presentation (requires python-pptx)."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            print("python-pptx is required for PPTX export. Install with: pip install python-pptx")
            return False
        
        import cv2
        import tempfile
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        total_steps = len(self.tutorial.steps)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add title text
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = self.tutorial.title
        p.font.size = Pt(44)
        p.font.bold = True
        
        for i, step in enumerate(self.tutorial.steps):
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Get image
            img = None
            if step.image_path and os.path.exists(step.image_path):
                img = cv2.imread(step.image_path)
            elif self.tutorial.video_path and os.path.exists(self.tutorial.video_path):
                cap = cv2.VideoCapture(self.tutorial.video_path)
                fps = cap.get(cv2.CAP_PROP_FPS) or 24.0
                cap.set(cv2.CAP_PROP_POS_FRAMES, int(step.timestamp * fps))
                ret, img = cap.read()
                cap.release()
            
            if img is not None:
                img = self._draw_overlay(img, step, i + 1)
                
                # Save temp image (don't delete yet - Windows locks it)
                tmp_path = os.path.join(tempfile.gettempdir(), f"tutomake_step_{i}.png")
                cv2.imwrite(tmp_path, img)
                slide.shapes.add_picture(tmp_path, Inches(0.5), Inches(0.5), height=Inches(6))
            
            # Step description
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"Step {i+1}: {step.description}"
            p.font.size = Pt(18)
            p.font.bold = True
            
            if step.instruction:
                instrBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.3))
                instrTf = instrBox.text_frame
                instrP = instrTf.paragraphs[0]
                instrP.text = step.instruction
                instrP.font.size = Pt(14)
                instrP.font.italic = True
            
            if self.progress_callback:
                self.progress_callback(int((i + 1) / total_steps * 100))
        
        prs.save(output_path)
        
        # Clean up temp files after save
        for i in range(len(self.tutorial.steps)):
            tmp_path = os.path.join(tempfile.gettempdir(), f"tutomake_step_{i}.png")
            try:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
            except:
                pass
        
        print(f"Exported PPTX to: {output_path}")
        return True
